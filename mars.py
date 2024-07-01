# M.A.R.S (Multi-modal AI Research System) by Sahil Bhoite
import streamlit as st
import os
import io
import pandas as pd
from PyPDF2 import PdfReader
from pptx import Presentation
import docx
import xlrd
from bs4 import BeautifulSoup
from langchain_text_splitters import RecursiveCharacterTextSplitter
from langchain_community.embeddings import GooglePalmEmbeddings
from langchain_community.llms import GooglePalm
from langchain_community.vectorstores import FAISS
from langchain.chains import ConversationalRetrievalChain
from langchain.memory import ConversationBufferMemory
from multiprocessing import Pool
import cv2
import pytesseract
import numpy as np
from langchain_core.messages import HumanMessage, AIMessage
import logging
from typing import Union
import speech_recognition as sr
from datetime import datetime
import spacy  # Added for Named Entity Recognition (NER)
from config import GOOGLE_API_KEY
import zipfile
import rarfile
import markdown
from textract import process
from transformers import pipeline
from langchain_community.chat_models import ChatOllama 

os.environ['GOOGLE_API_KEY'] = GOOGLE_API_KEY # Added for Google API Key for Google Palm 

nlp = spacy.load("en_core_web_sm") # Added for Named Entity Recognition (NER) 

def extract_text(file): # Added for various file types text extraction
    text = ""
    file_extension = file.name.split(".")[-1].lower()
    if file_extension == "pdf":
        text = extract_text_from_pdf(file)
    elif file_extension == "pptx":
        text = extract_text_from_ppt(file)
    elif file_extension == "py":
        text = extract_text_from_py(file)
    elif file_extension in ["doc", "docx"]:
        text = extract_text_from_docx(file)
    elif file_extension in ["xls", "xlsx"]:
        text = extract_text_from_excel(file)
    elif file_extension == "csv":
        text = extract_text_from_csv(file)
    elif file_extension == "html":
        text = extract_text_from_html(file)
    elif file_extension == "css":
        text = extract_text_from_css(file)
    elif file_extension == "json":
        text = extract_text_from_json(file)
    elif file_extension == "sql":
        text = extract_text_from_sql(file)
    elif file_extension == "txt":
        text = extract_text_from_txt(file)
    elif file_extension == "java":
        text = extract_text_from_java(file)
    elif file_extension in ["c", "h"]:
        text = extract_text_from_c(file)
    elif file_extension == "cpp":
        text = extract_text_from_cpp(file)
    elif file_extension == "js":
        text = extract_text_from_javascript(file)
    elif file_extension == "swift":
        text = extract_text_from_swift(file)
    elif file_extension == "r":
        text = extract_text_from_r(file)
    elif file_extension == "rs":
        text = extract_text_from_rust(file)
    elif file_extension in ["jpg", "jpeg", "png", "bmp"]:
        text = extract_text_from_image(file)
    elif file_extension == "xml":
        text = extract_text_from_xml(file)
    elif file_extension == "md":
        text = extract_text_from_md(file)
    elif file_extension == "tex":
        text = extract_text_from_tex(file)
    elif file_extension == "zip":
        text = extract_text_from_zip(file)
    elif file_extension == "rar":
        text = extract_text_from_rar(file)
    return text

# Function to extract text from a XML file
def extract_text_from_xml(xml_file):
    text = ""
    try:
        soup = BeautifulSoup(xml_file, 'xml')
        text = soup.get_text()
    except Exception as e:
        handle_file_processing_error("XML", e)
    return text

# Function to extract text from a PDF file
def extract_text_from_pdf(pdf_file):
    text = ""
    try:
        pdf_reader = PdfReader(pdf_file)
        for page in pdf_reader.pages:
            text += page.extract_text()
    except Exception as e:
        handle_file_processing_error("PDF", e)
    return text

# Function to extract text from a PowerPoint file
def extract_text_from_ppt(ppt_file):
    text = ""
    try:
        presentation = Presentation(ppt_file)
        for slide in presentation.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text
    except Exception as e:
        handle_file_processing_error("PPT", e)
    return text

# Function to extract text from a Python file
def extract_text_from_py(py_file):
    text = ""
    try:
        text = py_file.read().decode("utf-8")
    except Exception as e:
        handle_file_processing_error("Python", e)
    return text

# Function to extract text from a Word file
def extract_text_from_docx(docx_file):
    text = ""
    try:
        doc = docx.Document(docx_file)
        for paragraph in doc.paragraphs:
            text += paragraph.text + "\n"
    except Exception as e:
        handle_file_processing_error("DOCX", e)
    return text

# Function to extract text from a Excel file
def extract_text_from_excel(excel_file):
    text = ""
    try:
        workbook = xlrd.open_workbook(file_contents=excel_file.read())
        for sheet in workbook.sheets():
            for row in range(sheet.nrows):
                for col in range(sheet.ncols):
                    cell_value = sheet.cell(row, col).value
                    if isinstance(cell_value, str):
                        text += cell_value + " "
    except Exception as e:
        handle_file_processing_error("Excel", e)
    return text

# Function to extract text from a CSV file
def extract_text_from_csv(csv_file):
    text = ""
    try:
        csv_data = io.BytesIO(csv_file.read())
        df = pd.read_csv(csv_data)
        text = df.to_string(index=False)
    except Exception as e:
        handle_file_processing_error("CSV", e)
    return text

# Function to extract text from a HTML file
def extract_text_from_html(html_file):
    text = ""
    try:
        soup = BeautifulSoup(html_file, 'html.parser')
        text = soup.get_text()
    except Exception as e:
        handle_file_processing_error("HTML", e)
    return text

# Function to extract text from a CSS file
def extract_text_from_css(css_file):
    text = ""
    try:
        text = css_file.read().decode("utf-8")
    except Exception as e:
        handle_file_processing_error("CSS", e)
    return text

# Function to extract text from a Wav file
def extract_text_from_audio(audio_file):
    text = ""
    try:
        r = sr.Recognizer()
        with sr.AudioFile(audio_file) as source:
            audio_data = r.record(source)
            text = r.recognize_google(audio_data)
    except Exception as e:
        handle_file_processing_error("audio", e)
    return text

# Function to extract text from a JSON file
def extract_text_from_json(json_file):
    text = ""
    try:
        with open(json_file.name, "r") as f:
            text = f.read()
    except Exception as e:
        handle_file_processing_error("JSON", e)
    return text

# Function to extract text from a SQL file
def extract_text_from_sql(sql_file):
    text = ""
    try:
        with open(sql_file.name, "r") as f:
            text = f.read()
    except Exception as e:
        handle_file_processing_error("SQL", e)
    return text

# Function to extract text from a TXT file
def extract_text_from_txt(txt_file):
    text = ""
    try:
        with open(txt_file.name, "r") as f:
            text = f.read()
    except Exception as e:
        handle_file_processing_error("TXT", e)
    return text

# Function to extract text from a Java file
def extract_text_from_java(java_file):
    text = ""
    try:
        text = java_file.read().decode("utf-8")
    except Exception as e:
        handle_file_processing_error("Java", e)
    return text

# Function to extract text from a C file
def extract_text_from_c(c_file):
    text = ""
    try:
        text = c_file.read().decode("utf-8")
    except Exception as e:
        handle_file_processing_error("C", e)
    return text

# Function to extract text from a C++ file
def extract_text_from_cpp(cpp_file):
    text = ""
    try:
        text = cpp_file.read().decode("utf-8")
    except Exception as e:
        handle_file_processing_error("C++", e)
    return text

# Function to extract text from a Javascript file
def extract_text_from_javascript(js_file):
    text = ""
    try:
        text = js_file.read().decode("utf-8")
    except Exception as e:
        handle_file_processing_error("JavaScript", e)
    return text

# Function to extract text from a Swift file
def extract_text_from_swift(swift_file):
    text = ""
    try:
        text = swift_file.read().decode("utf-8")
    except Exception as e:
        handle_file_processing_error("Swift", e)
    return text

# Function to extract text from a R file
def extract_text_from_r(r_file):
    text = ""
    try:
        text = r_file.read().decode("utf-8")
    except Exception as e:
        handle_file_processing_error("R", e)
    return text

# Function to extract text from a Rust file
def extract_text_from_rust(rs_file):
    text = ""
    try:
        text = rs_file.read().decode("utf-8")
    except Exception as e:
        handle_file_processing_error("Rust", e)
    return text

# Function to extract text from a Images
def extract_text_from_image(image_file):
    text = ""
    try:
        image = cv2.imdecode(np.frombuffer(image_file.read(), np.uint8), cv2.IMREAD_COLOR)
        gray = cv2.cvtColor(image, cv2.COLOR_BGR2GRAY)
        text = pytesseract.image_to_string(gray)
    except Exception as e:
        handle_file_processing_error("image", e)
    return text

# Function to extract text from a Markdown file
def extract_text_from_md(md_file):
    text = ""
    try:
        text = md_file.read().decode("utf-8")
        text = markdown.markdown(text)
    except Exception as e:
        handle_file_processing_error("Markdown", e)
    return text

# Function to extract text from a LaTeX file
def extract_text_from_tex(tex_file):
    text = ""
    try:
        text = process(tex_file.name).decode("utf-8")
    except Exception as e:
        handle_file_processing_error("LaTeX", e)
    return text

# Function to extract text from a Zip file
def extract_text_from_zip(zip_file):
    text = ""
    try:
        with zipfile.ZipFile(zip_file, 'r') as z:
            for filename in z.namelist():
                with z.open(filename) as f:
                    text += extract_text(io.BytesIO(f.read()))
    except Exception as e:
        handle_file_processing_error("ZIP", e)
    return text

# Function to extract text from a Rar file
def extract_text_from_rar(rar_file):
    text = ""
    try:
        with rarfile.RarFile(rar_file, 'r') as rf:
            for filename in rf.namelist():
                with rf.open(filename) as f:
                    text += extract_text(io.BytesIO(f.read()))
    except Exception as e:
        handle_file_processing_error("RAR", e)
    return text

# Function to handle file processing errors
def handle_file_processing_error(file_type: str, error: Exception):
    st.error(f"Error processing {file_type} file: {error}")
    logger.exception(f"Error processing {file_type} file", exc_info=True)

# Function to handle AI model interaction errors
def handle_model_interaction_error(error: Exception):
    st.error(f"Error interacting with AI model: {error}")
    logger.exception("Error interacting with AI model", exc_info=True)

# Create a logger
logger = logging.getLogger(__name__)

# Function to handle input validation errors
def validate_user_input(user_input: Union[st.file_uploader, str]):
    if not user_input:
        st.warning("Please provide valid input.")
        return False
    return True

# Function to handle user feedback
def handle_user_feedback(feedback: str):
    st.success("Thank you for your feedback!")

# Function to log important events and interactions
def log_event(event: str):
    logger.info(event)

def main():
    st.set_page_config(page_title="M.A.R.S 🚀", layout="wide")
    st.header("Multi-model AI Research System")
    user_question = st.chat_input("Ask Questions about Everything")

    if "conversation" not in st.session_state or not st.session_state.conversation:
        st.session_state.conversation = None
        st.session_state.chat_history = []
        st.session_state.files_uploaded = False

    with st.sidebar:
        st.title("M.A.R.S")
        model_mode = st.toggle("Online Mode")

        st.subheader("Upload your Files here")
        files = st.file_uploader("Upload your Files and Click on the NEXT Button", accept_multiple_files=True)

        if st.button("NEXT"):
            if validate_user_input(files):
                with st.spinner("Processing your Files..."):
                    raw_text = ""
                    for file in files:
                        try:
                            if file.type.startswith('audio'):
                                raw_text += extract_text_from_audio(file)
                            else:
                                raw_text += extract_text(file)
                        except Exception as e:
                            handle_file_processing_error(file.name.split(".")[-1].lower(), e)

                    text_chunks = get_text_chunks(raw_text)
                    vector_store = get_vector_store(text_chunks)

                    if model_mode:
                        st.session_state.conversation = get_conversational_chain_online(vector_store)
                    else:
                        st.session_state.conversation = get_conversational_chain_offline(vector_store)

                    st.session_state.files_uploaded = True
                    st.success("Processing Done!")
            else:
                st.warning("Please upload at least one file.")

    if user_question:
        user_input(user_question)

    if not st.session_state.files_uploaded:
        st.warning("Start the chat by uploading your files.")
    elif st.session_state.files_uploaded and not files:
        st.session_state.files_uploaded = False



# Function to split text into chunks
def get_text_chunks(text):
    text_splitter = RecursiveCharacterTextSplitter(chunk_size=1000, chunk_overlap=200)
    chunks = text_splitter.split_text(text)
    return chunks
    
# Function to create a vector store
def get_vector_store(text_chunks):
    embeddings = GooglePalmEmbeddings()
    vector_store = FAISS.from_texts(text_chunks, embedding=embeddings)
    return vector_store

# Function to create a conversational chain using Sol
def get_conversational_chain_offline(vector_store):
    sol_model = ChatOllama(model="Sol")  # Initialize the ChatOllama instance with the Sol model (Made for MARS)
    memory = ConversationBufferMemory(memory_key="chat_history", return_messages=True)
    conversation_chain = ConversationalRetrievalChain.from_llm(llm=sol_model, retriever=vector_store.as_retriever(), memory=memory)
    return conversation_chain

# Function to create a conversational chain using Palm
def get_conversational_chain_online(vector_store):
    llm = GooglePalm()
    memory = ConversationBufferMemory(memory_key="chat_history", return_messages=True)
    conversation_chain = ConversationalRetrievalChain.from_llm(llm=llm, retriever=vector_store.as_retriever(), memory=memory)
    return conversation_chain

def user_input(user_question):
    if st.session_state.conversation:
        try:
            # Check if the user's message contains a feedback mention
            is_feedback = "@mars" in user_question.lower()
            is_feedback = is_feedback or "@mars" in user_question.upper()
            
            # If the message contains a feedback mention, handle it
            if is_feedback:
                handle_user_feedback(user_question)
            else:
                # Otherwise, proceed with the conversation
                response = st.session_state.conversation({'question': user_question})
                if 'chat_history' in response:
                    st.session_state.chat_history = response['chat_history']
                if st.session_state.chat_history:
                    for message in st.session_state.chat_history:
                        if isinstance(message, HumanMessage):
                            with st.chat_message("User"):
                                st.write(message.content)
                        elif isinstance(message, AIMessage):
                            with st.chat_message("AI"):
                                st.write(message.content)
                    
                    # Auto-scroll to the end of the chat with smooth behavior
                    st.markdown(
                        """
                        <script>
                        var element = document.getElementById("end-of-chat");
                        element.scrollIntoView({behavior: "smooth"});
                        </script>
                        """,
                        unsafe_allow_html=True,
                    )
        except Exception as e:
            handle_model_interaction_error(e)
            st.error("An error occurred during conversation. Please try again.")
    else:
        st.warning("Please upload files and click 'NEXT' to start the conversation.")
# Running the main function
if __name__ == "__main__":
    main()


